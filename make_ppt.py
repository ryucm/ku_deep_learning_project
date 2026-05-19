"""
패션 시대 분류 및 코디 추천 프로젝트 발표 자료
python make_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import os

BASE   = Path(__file__).parent
OUTPUT = BASE / "presentation.pptx"

SW, SH = 13.33, 7.5

C_DARK  = RGBColor(0x0d, 0x11, 0x17)
C_NAVY  = RGBColor(0x16, 0x21, 0x3e)
C_BLUE  = RGBColor(0x1a, 0x72, 0xa8)
C_LBLUE = RGBColor(0xe8, 0xf2, 0xfc)
C_WHITE = RGBColor(0xff, 0xff, 0xff)
C_LGRAY = RGBColor(0xf2, 0xf4, 0xf7)
C_MGRAY = RGBColor(0xd0, 0xd8, 0xe8)
C_TEXT  = RGBColor(0x22, 0x22, 0x33)
C_GRAY  = RGBColor(0x77, 0x77, 0x88)
C_BOX   = RGBColor(0xf8, 0xfa, 0xfc)

KO = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, lw=1.0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=11, bold=False, color=C_TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = KO
    return tb

def hdr(slide, section, title):
    set_bg(slide, C_LGRAY)
    box(slide, 0, 0, SW, 1.05, fill=C_NAVY)
    txt(slide, section, 0.45, 0.08, 10, 0.36, size=10,
        color=RGBColor(0x99, 0xbb, 0xdd))
    txt(slide, title, 0.45, 0.44, 10, 0.54, size=20, bold=True, color=C_WHITE)

def pic(slide, path, l, t, w=None, h=None):
    if not os.path.exists(str(path)):
        return
    kw = {'left': Inches(l), 'top': Inches(t)}
    if w: kw['width']  = Inches(w)
    if h: kw['height'] = Inches(h)
    slide.shapes.add_picture(str(path), **kw)

def tbl_cell(cell, text, size=9, bold=False, color=C_TEXT,
             bg_color=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = KO
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    tf.margin_left  = Inches(0.08)
    tf.margin_top   = Inches(0.05)


# ── Slide 1: 타이틀 ──────────────────────────────
def slide_title():
    sl = new_slide()
    set_bg(sl, C_DARK)
    box(sl, 0, 0, 0.5, SH, fill=C_NAVY)
    txt(sl, "패션 시대 분류 및\n코디 추천 프로젝트",
        1.0, 1.9, 8.5, 2.4, size=34, bold=True, color=C_WHITE)
    txt(sl, "Fashion Era Classification & Outfit Recommendation",
        1.0, 4.45, 9.5, 0.55, size=14,
        color=RGBColor(0xaa, 0xcc, 0xee), italic=True)
    box(sl, 1.0, 5.15, 6.0, 0.04, fill=C_BLUE)
    txt(sl, "류창민  |  딥러닝 프로젝트",
        1.0, 5.3, 6.0, 0.5, size=12, color=C_GRAY)
    txt(sl, "AI Hub 연도별 패션 데이터  |  PyTorch 2.5.1  |  Apple M4 Pro MPS",
        1.0, 5.8, 9.0, 0.45, size=10, color=C_GRAY)


# ── Slide 2: Contents ────────────────────────────
def slide_contents():
    sl = new_slide()
    set_bg(sl, C_WHITE)
    box(sl, 0, 0, SW, 0.85, fill=C_NAVY)
    txt(sl, "Contents", 0.5, 0.08, 5, 0.6, size=24, bold=True, color=C_WHITE)
    txt(sl, "패션 시대 분류 및 코디 추천 프로젝트",
        0.5, 0.52, 9, 0.35, size=11,
        color=RGBColor(0xaa, 0xbb, 0xcc))

    sections = [
        ("01", "프로젝트 소개",       ["• 프로젝트 목표", "• 시스템 구조"]),
        ("02", "데이터 파이프라인",   ["• 학습 데이터 준비", "• 데이터셋 구성 및 전처리"]),
        ("03", "모델 구축 및 학습",   ["• 모델 설계", "• 학습 및 검증"]),
        ("04", "마무리",              ["• 프로젝트 결과", "• 향후 발전 방향"]),
    ]
    bw = 2.9
    for i, (num, title, bullets) in enumerate(sections):
        lx = 0.35 + i * (bw + 0.18)
        box(sl, lx, 1.0, bw, 5.7, fill=C_BOX, line=C_MGRAY)
        txt(sl, num, lx, 1.1, bw, 1.4,
            size=52, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        txt(sl, title, lx, 2.55, bw, 0.55,
            size=13, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
        txt(sl, "\n".join(bullets), lx + 0.18, 3.2, bw - 0.35, 2.3,
            size=11, color=C_GRAY)


# ── Slide 3: I-01 프로젝트 목표 ──────────────────
def slide_goal():
    sl = new_slide()
    hdr(sl, "I. 프로젝트 소개", "01. 프로젝트 목표")
    txt(sl,
        "사용자의 옷 이미지를 입력받아 스타일 벡터를 추출하고\n"
        "패션 데이터셋과의 유사도를 기반으로 어울리는 코디를 추천합니다.",
        0.45, 1.2, 12.0, 0.9, size=12)

    goals = [
        ("시대 분류 학습",
         "연도별 패션 이미지로\nCNN을 학습하여\n1990 / 2000 / 2010 / 2019\n스타일 패턴 파악"),
        ("스타일 벡터 추출",
         "학습된 CNN의 encode()\n메서드로 이미지를\n512-dim 스타일 벡터로\n변환"),
        ("코디 추천",
         "입력 이미지 벡터와\n데이터셋 전체 벡터의\n코사인 유사도 계산\n→ Top-K 코디 반환"),
    ]
    bw = 3.7
    for i, (title, body) in enumerate(goals):
        lx = 0.45 + i * (bw + 0.45)
        box(sl, lx, 2.2, bw, 4.0, fill=C_WHITE, line=C_NAVY, lw=1.5)
        box(sl, lx, 2.2, bw, 0.62, fill=C_NAVY)
        txt(sl, title, lx, 2.22, bw, 0.58,
            size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, body, lx + 0.18, 3.0, bw - 0.35, 3.0,
            size=11, color=C_TEXT, align=PP_ALIGN.CENTER)

    box(sl, 0, 6.78, SW, 0.58, fill=C_NAVY)
    txt(sl,
        "단순한 분류를 넘어 스타일 유사도 기반의 실용적인 코디 추천 시스템 구현",
        0.3, 6.82, SW - 0.6, 0.48,
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ── Slide 4: I-02 시스템 구조 ────────────────────
def slide_system():
    sl = new_slide()
    hdr(sl, "I. 프로젝트 소개", "02. 시스템 구조")
    txt(sl, "데이터 전처리, 모델 학습, 특징 추출, 추천이 각각 독립적인 모듈로 구성됩니다.",
        0.45, 1.15, 12.0, 0.45, size=12)

    modules = [
        ("01", "데이터 준비\n& 전처리",  "3000×4000 → 256×256\n오프라인 리사이즈\n데이터 증강"),
        ("02", "모델 학습",              "CustomCNN / ResNet18\n2단계 Transfer Learning\n체크포인트 자동 저장"),
        ("03", "특징 벡터 추출",         "encode() → 512-dim\n21,985장 전체 인코딩\nfeatures.npy 저장"),
        ("04", "코디 추천",              "코사인 유사도 계산\nTop-K 검색\n결과 시각화"),
    ]
    bw = 2.7
    for i, (num, title, body) in enumerate(modules):
        lx = 0.4 + i * (bw + 0.5)
        box(sl, lx, 1.8, bw, 4.0, fill=C_BOX, line=C_MGRAY)
        box(sl, lx, 1.8, bw, 0.58, fill=C_BLUE)
        txt(sl, num, lx, 1.82, bw, 0.54,
            size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, title, lx + 0.1, 2.48, bw - 0.2, 0.7,
            size=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        txt(sl, body, lx + 0.15, 3.3, bw - 0.3, 2.3,
            size=10, color=C_TEXT, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(sl, "▶", lx + bw + 0.08, 3.55, 0.35, 0.45,
                size=16, color=C_NAVY, align=PP_ALIGN.CENTER)

    box(sl, 0.4, 6.1, SW - 0.8, 0.52, fill=C_MGRAY)
    txt(sl, "모델 저장 및 추적   |   best_resnet18.pth 자동 저장, Val Accuracy 모니터링",
        0.5, 6.14, SW - 1.0, 0.42, size=11, color=C_NAVY, align=PP_ALIGN.CENTER)


# ── Slide 5: II-01 데이터 준비 ───────────────────
def slide_data():
    sl = new_slide()
    hdr(sl, "II. 데이터 파이프라인", "01. 패션 데이터 준비")
    txt(sl, "출처: AI Hub — 연도별 패션 선호도 파악 및 추천 데이터 (CC BY-SA)",
        0.45, 1.2, 10.0, 0.38, size=12)

    box(sl, 0.4, 1.65, 2.8, 0.38, fill=C_BLUE)
    txt(sl, "데이터 준비 : AI Hub 패션 데이터셋",
        0.5, 1.67, 2.7, 0.34, size=11, bold=True, color=C_WHITE)

    tbl_data = [
        ["항목",     "내용"],
        ["이미지 수",  "Training 14,438장  /  Validation 7,547장  (총 21,985장)"],
        ["원본 해상도", "3000 × 4000 px (패션 전신 사진)"],
        ["레이블",    "촬영 연도 4클래스 (1990 / 2000 / 2010 / 2019)"],
        ["성별",      "남성 / 여성"],
        ["부가 정보",  "스타일명 27종 (normcore, hiphop, metrosexual 등), 사용자 설문 응답"],
    ]
    ts = sl.shapes.add_table(
        len(tbl_data), 2,
        Inches(0.4), Inches(2.1),
        Inches(8.6), Inches(3.6)
    )
    t = ts.table
    t.columns[0].width = Inches(2.2)
    t.columns[1].width = Inches(6.4)
    for ri, row in enumerate(tbl_data):
        for ci, val in enumerate(row):
            is_h = ri == 0
            tbl_cell(t.cell(ri, ci), val,
                     bold=is_h, size=10 if not is_h else 11,
                     color=C_WHITE if is_h else C_TEXT,
                     bg_color=C_NAVY if is_h else (C_BOX if ri % 2 == 0 else C_WHITE),
                     align=PP_ALIGN.CENTER if (is_h or ci == 0) else PP_ALIGN.LEFT)

    txt(sl, "데이터 특성", 9.3, 1.67, 3.7, 0.38, size=12, bold=True, color=C_BLUE)
    chars = [
        "• 단순 카테고리가 아닌\n  시대별 스타일 변화 데이터",
        "• 인접 연도(2010↔2019)는\n  시각적 차이가 작아 난이도 높음",
        "• 원본 3000×4000px →\n  256×256 사전 리사이즈 처리",
    ]
    for j, c in enumerate(chars):
        box(sl, 9.3, 2.15 + j * 1.3, 3.7, 1.15, fill=C_LBLUE, line=C_BLUE, lw=0.8)
        txt(sl, c, 9.42, 2.22 + j * 1.3, 3.5, 1.0, size=10)


# ── Slide 6: II-02 전처리 ────────────────────────
def slide_preprocess():
    sl = new_slide()
    hdr(sl, "II. 데이터 파이프라인", "02. 데이터셋 구성 및 전처리")

    # 왼쪽: 오프라인 리사이즈
    box(sl, 0.4, 1.15, 5.9, 0.44, fill=C_NAVY)
    txt(sl, "1. 오프라인 사전 리사이즈",
        0.5, 1.17, 5.7, 0.38, size=13, bold=True, color=C_WHITE)

    items = [
        ("원본 이미지\n3000×4000px\n79 GB",     0.55),
        ("256×256\n사전 리사이즈\n411 MB",       2.6),
        ("학습 시\nCenterCrop\n224×224",         4.65),
    ]
    for text, lx in items:
        box(sl, lx, 1.75, 1.7, 1.35, fill=C_BOX, line=C_MGRAY)
        txt(sl, text, lx + 0.05, 1.95, 1.6, 0.95,
            size=10, color=C_TEXT, align=PP_ALIGN.CENTER)
    for ax in [2.35, 4.38]:
        txt(sl, "→", ax, 2.3, 0.35, 0.5,
            size=20, color=C_NAVY, align=PP_ALIGN.CENTER)

    txt(sl, "• 학습 속도 5배 향상 (3.54 s/batch → 0.69 s/batch)\n"
            "• 79 GB 원본 → 411 MB 리사이즈 데이터 (98.3% 절감)",
        0.55, 3.22, 5.6, 0.85, size=10)

    # 오른쪽: 데이터 증강
    box(sl, 6.8, 1.15, 6.1, 0.44, fill=C_NAVY)
    txt(sl, "2. 학습 데이터 증강 (Data Augmentation)",
        6.9, 1.17, 5.9, 0.38, size=13, bold=True, color=C_WHITE)

    augs = [
        ("RandomHorizontalFlip",   "좌우 반전  (p = 0.5)"),
        ("RandomCrop(224, pad=16)", "랜덤 크롭 (패딩 후 자름)"),
        ("ColorJitter",            "밝기 / 대비 / 채도 무작위 변환"),
        ("Normalize",              "ImageNet mean / std 정규화"),
    ]
    for j, (name, desc) in enumerate(augs):
        box(sl, 6.8, 1.72 + j * 1.0, 6.1, 0.85, fill=C_BOX, line=C_MGRAY)
        txt(sl, name, 7.0, 1.79 + j * 1.0, 3.5, 0.36, size=11, bold=True, color=C_BLUE)
        txt(sl, desc, 7.0, 2.14 + j * 1.0, 5.7, 0.34, size=10)
    txt(sl, "※ Validation / Test: CenterCrop(224) + Normalize만 적용",
        6.8, 5.85, 6.0, 0.38, size=9, color=C_GRAY, italic=True)


# ── Slide 7: III-01 모델 설계 (아키텍처) ─────────
def slide_model_arch():
    sl = new_slide()
    hdr(sl, "III. 모델 구축 및 학습", "01. 모델 설계")
    box(sl, 0.4, 1.12, 3.5, 0.42, fill=C_BLUE)
    txt(sl, "2. 모델 정의 (아키텍처)",
        0.5, 1.15, 3.3, 0.36, size=12, bold=True, color=C_WHITE)

    # CustomCNN
    box(sl, 0.4, 1.68, 5.8, 5.15, fill=C_WHITE, line=C_MGRAY)
    txt(sl, "CustomCNN", 0.5, 1.73, 5.6, 0.52,
        size=16, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    txt(sl, "From Scratch — 구조 이해 목적",
        0.5, 2.25, 5.6, 0.38, size=10, color=C_GRAY,
        align=PP_ALIGN.CENTER, italic=True)

    cnn = [
        ("입력  (3 × 224 × 224)",              False, C_TEXT),
        ("ConvBlock × 4",                       True,  C_NAVY),
        ("(Conv → BN → ReLU → MaxPool)",        False, C_GRAY),
        ("AdaptiveAvgPool → Flatten",            False, C_TEXT),
        ("FC(256→512) → ReLU → Dropout(0.5)",   False, C_TEXT),
        ("encode()  →  512-dim 벡터",            True,  C_BLUE),
        ("FC(512→4)  →  클래스 확률",            True,  C_NAVY),
    ]
    for j, (s, b, c) in enumerate(cnn):
        txt(sl, s, 0.6, 2.72 + j * 0.56, 5.4, 0.5,
            size=10, bold=b, color=c, align=PP_ALIGN.CENTER)

    # ResNet18
    box(sl, 6.9, 1.68, 5.9, 5.15, fill=C_WHITE, line=C_MGRAY)
    txt(sl, "ResNet18", 7.0, 1.73, 5.7, 0.52,
        size=16, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    txt(sl, "Transfer Learning — 성능 향상 목적",
        7.0, 2.25, 5.7, 0.38, size=10, color=C_GRAY,
        align=PP_ALIGN.CENTER, italic=True)

    rn = [
        ("입력  (3 × 224 × 224)",               False, C_TEXT),
        ("backbone",                             True,  C_NAVY),
        ("(ImageNet 사전학습 ResNet18",           False, C_GRAY),
        (" FC 레이어 제거)",                      False, C_GRAY),
        ("Flatten  →  512-dim",                  False, C_TEXT),
        ("encode()  →  512-dim 벡터",             True,  C_BLUE),
        ("FC(512→4)  →  클래스 확률",             True,  C_NAVY),
    ]
    for j, (s, b, c) in enumerate(rn):
        txt(sl, s, 7.1, 2.72 + j * 0.56, 5.6, 0.5,
            size=10, bold=b, color=c, align=PP_ALIGN.CENTER)

    box(sl, 0.4, 6.66, SW - 0.8, 0.58, fill=C_LBLUE, line=C_BLUE, lw=0.8)
    txt(sl,
        "두 모델 모두 encode() 메서드로 512-dim 특징 벡터 추출 → 추천 시스템과 동일한 인터페이스",
        0.6, 6.72, SW - 1.2, 0.46, size=11, color=C_NAVY, align=PP_ALIGN.CENTER)


# ── Slide 8: III-01 Transfer Learning ────────────
def slide_transfer():
    sl = new_slide()
    hdr(sl, "III. 모델 구축 및 학습", "01. 모델 설계")
    box(sl, 0.4, 1.12, 4.2, 0.42, fill=C_BLUE)
    txt(sl, "ResNet18 — 2단계 Transfer Learning",
        0.5, 1.15, 4.1, 0.36, size=12, bold=True, color=C_WHITE)
    txt(sl,
        "backbone 과적합 방지를 위해 2단계로 나눠 학습률을 제어합니다.",
        0.45, 1.65, 12.5, 0.38, size=11)

    # Stage 1
    box(sl, 0.45, 2.15, 5.8, 2.55, fill=C_BOX, line=C_MGRAY)
    box(sl, 0.45, 2.15, 5.8, 0.52, fill=C_NAVY)
    txt(sl, "Stage 1  (5 epoch) — backbone 동결",
        0.55, 2.17, 5.6, 0.46, size=13, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    s1 = ["• ImageNet 사전학습 가중치 보호",
          "• backbone.requires_grad = False",
          "• FC 레이어(head)만 업데이트",
          "• head lr = 1e-3"]
    for j, l in enumerate(s1):
        txt(sl, l, 0.65, 2.78 + j * 0.46, 5.4, 0.42, size=11)

    txt(sl, "▶", 6.38, 3.3, 0.6, 0.5, size=22, color=C_NAVY, align=PP_ALIGN.CENTER)

    # Stage 2
    box(sl, 7.1, 2.15, 5.8, 2.55, fill=C_BOX, line=C_MGRAY)
    box(sl, 7.1, 2.15, 5.8, 0.52, fill=C_NAVY)
    txt(sl, "Stage 2  (25 epoch) — 전체 fine-tuning",
        7.2, 2.17, 5.6, 0.46, size=13, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    s2 = ["• backbone.requires_grad = True",
          "• backbone lr = 1e-4  (사전학습 보존)",
          "• head lr = 1e-3  (적극 학습)",
          "• 패션 도메인 적응"]
    for j, l in enumerate(s2):
        txt(sl, l, 7.25, 2.78 + j * 0.46, 5.4, 0.42, size=11)

    # 공통 설정
    box(sl, 0.45, 4.95, SW - 0.9, 0.48, fill=C_NAVY)
    txt(sl, "공통 설정",
        0.6, 4.97, 2.0, 0.42, size=12, bold=True, color=C_WHITE)
    commons = [
        ("Loss",           "CrossEntropy\n+ Label Smoothing 0.1"),
        ("Optimizer",      "Adam\n(weight_decay = 1e-4)"),
        ("Scheduler",      "CosineAnnealingLR"),
        ("Early Stopping", "patience = 10"),
    ]
    bw4 = (SW - 1.0) / 4
    for i, (k, v) in enumerate(commons):
        lx = 0.45 + i * (bw4 + 0.05)
        box(sl, lx, 5.52, bw4, 1.65, fill=C_WHITE, line=C_MGRAY)
        txt(sl, k, lx, 5.58, bw4, 0.42,
            size=12, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
        txt(sl, v, lx + 0.1, 6.08, bw4 - 0.2, 0.95,
            size=10, color=C_TEXT, align=PP_ALIGN.CENTER)


# ── Slide 9: III-02 학습 결과 ────────────────────
def slide_training():
    sl = new_slide()
    hdr(sl, "III. 모델 구축 및 학습", "02. 학습 및 검증")
    box(sl, 0.4, 1.12, 2.0, 0.42, fill=C_BLUE)
    txt(sl, "3. 훈련 및 검증",
        0.5, 1.15, 1.9, 0.36, size=12, bold=True, color=C_WHITE)

    funcs = [
        ("train_epoch",  "모델을 한 epoch 동안 학습\n역전파로 파라미터 업데이트"),
        ("validate",     "검증 세트 평가\nLoss & Accuracy 측정"),
        ("run_loop",     "N epoch 반복 실행\nEarly Stopping 포함\n최적 모델 자동 저장"),
    ]
    bw3 = 3.6
    for i, (name, desc) in enumerate(funcs):
        lx = 0.4 + i * (bw3 + 0.37)
        box(sl, lx, 1.68, bw3, 2.35, fill=C_BOX, line=C_MGRAY)
        box(sl, lx, 1.68, bw3, 0.5, fill=C_NAVY)
        txt(sl, name, lx, 1.7, bw3, 0.46,
            size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, desc, lx + 0.15, 2.28, bw3 - 0.3, 1.55,
            size=10, color=C_TEXT, align=PP_ALIGN.CENTER)
        if i < 2:
            txt(sl, "▶", lx + bw3 + 0.03, 2.66, 0.33, 0.42,
                size=18, color=C_NAVY, align=PP_ALIGN.CENTER)

    txt(sl, "랜덤 기준선 (4클래스): 25.0%",
        0.45, 4.18, 8.0, 0.36, size=11, color=C_GRAY, italic=True)

    result_data = [
        ["모델",               "Best Val Acc", "Train Acc", "비고"],
        ["CustomCNN",          "37.4%",        "~44%",      "Epoch 9 최고, 이후 과적합"],
        ["ResNet18 (1단계)",   "42.5%",        "93.0%",     "심각한 과적합 (Train-Val 갭 50%)"],
        ["ResNet18 (2단계)",   "44.1%",        "~65%",      "차등 학습률 + Label Smoothing 적용"],
    ]
    ts = sl.shapes.add_table(4, 4,
                              Inches(0.4), Inches(4.62),
                              Inches(SW - 0.8), Inches(2.5))
    t = ts.table
    for ci, cw in enumerate([3.1, 2.3, 2.0, 4.8]):
        t.columns[ci].width = Inches(cw)
    for ri, row in enumerate(result_data):
        for ci, val in enumerate(row):
            is_h  = ri == 0
            is_b  = ri == 3
            tbl_cell(t.cell(ri, ci), val,
                     bold=is_h or is_b, size=10,
                     color=C_WHITE if is_h else (C_NAVY if is_b else C_TEXT),
                     bg_color=C_NAVY if is_h else (C_LBLUE if is_b else
                               (C_BOX if ri % 2 == 0 else C_WHITE)),
                     align=PP_ALIGN.CENTER if ci < 3 else PP_ALIGN.LEFT)


# ── Slide 10: III-02 혼동 행렬 ───────────────────
def slide_confusion():
    sl = new_slide()
    hdr(sl, "III. 모델 구축 및 학습", "02. 학습 결과 — 혼동 행렬")
    txt(sl, "ResNet18 (2단계 Transfer Learning)  |  Val Accuracy: 44.1%",
        0.45, 1.18, 12.0, 0.4, size=13, bold=True, color=C_NAVY)
    pic(sl, str(BASE / "confusion_matrix_resnet18.png"), 0.45, 1.72, 6.8)

    txt(sl, "결과 분석", 7.85, 1.72, 5.0, 0.38,
        size=13, bold=True, color=C_BLUE)
    analyses = [
        ("1990년대", "다른 시대와 혼동 가장 큼\nF1 = 0.36"),
        ("2000년대", "Recall 높음 (0.74) but Precision 낮음 (0.35)\n→ 과잉 예측 경향"),
        ("2010년대", "F1 = 0.48  (중간 성능)"),
        ("2019년대", "Precision 높음 (0.65) but Recall 낮음 (0.31)\n→ 과소 예측 경향"),
    ]
    for j, (year, note) in enumerate(analyses):
        box(sl, 7.85, 2.2 + j * 1.28, 5.0, 1.15, fill=C_BOX, line=C_MGRAY)
        txt(sl, year, 8.0, 2.26 + j * 1.28, 4.7, 0.36,
            size=12, bold=True, color=C_NAVY)
        txt(sl, note, 8.0, 2.6 + j * 1.28, 4.7, 0.7,
            size=10)


# ── Slide 11: IV-01 프로젝트 결과 ────────────────
def slide_result():
    sl = new_slide()
    hdr(sl, "IV. 결론", "01. 프로젝트 결과")
    txt(sl,
        "학습된 ResNet18으로 스타일 벡터를 추출하고, 코사인 유사도로 어울리는 코디를 추천합니다.",
        0.45, 1.18, 12.5, 0.4, size=12)

    pipeline = [
        ("옷 이미지\n입력",                "입력"),
        ("encode()\n→ 512-dim\n벡터 추출", "특징 추출"),
        ("코사인 유사도\n계산\n(29,532장)", "검색"),
        ("Top-5\n코디 추천",               "출력"),
    ]
    pw = 2.55
    for i, (body, label) in enumerate(pipeline):
        lx = 0.45 + i * (pw + 0.62)
        box(sl, lx, 1.75, pw, 1.55, fill=C_NAVY)
        txt(sl, body, lx, 1.78, pw, 1.49,
            size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, label, lx, 3.38, pw, 0.32,
            size=9, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(sl, "→", lx + pw + 0.1, 2.38, 0.45, 0.48,
                size=20, color=C_NAVY, align=PP_ALIGN.CENTER)

    pic(sl, str(BASE / "recommendation_result.png"), 0.45, 3.85, SW - 0.9)


# ── Slide 12: IV-02 향후 발전 방향 ───────────────
def slide_future():
    sl = new_slide()
    hdr(sl, "IV. 결론", "02. 향후 발전 방향")
    txt(sl, "패션 시대 분류 및 스타일 추천 시스템을 기반으로 다양한 서비스로 확장할 수 있습니다.",
        0.45, 1.18, 12.5, 0.4, size=12)

    items = [
        ("정확도 향상",   "ViT / EfficientNet 등\n최신 모델 적용"),
        ("멀티모달",      "텍스트 + 이미지\n결합 추천"),
        ("개인화",        "사용자 이력 기반\n선호도 학습"),
        ("트렌드 분석",   "SNS 크롤링으로\n최신 트렌드 반영"),
        ("실시간 서비스", "FastAPI 배포\n모바일 앱 연동"),
        ("체형 매칭",     "신체 정보 결합\n체형별 코디 추천"),
    ]
    iw, ih = 3.8, 2.05
    for i, (title, body) in enumerate(items):
        col = i % 3; row = i // 3
        lx = 0.45 + col * (iw + 0.4)
        ty = 1.82  + row * (ih + 0.38)
        box(sl, lx, ty, iw, ih, fill=C_BOX, line=C_MGRAY)
        txt(sl, title, lx, ty + 0.15, iw, 0.45,
            size=13, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
        txt(sl, body, lx + 0.1, ty + 0.72, iw - 0.2, 1.1,
            size=11, color=C_TEXT, align=PP_ALIGN.CENTER)


# ── Slide 13: 감사합니다 ──────────────────────────
def slide_end():
    sl = new_slide()
    set_bg(sl, C_DARK)
    box(sl, 0, 0, SW, 0.5, fill=C_NAVY)
    txt(sl, "감사합니다",
        0, SH / 2 - 0.55, SW, 1.1,
        size=46, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "패션 시대 분류 및 코디 추천 프로젝트",
        0, SH / 2 + 0.72, SW, 0.5,
        size=14, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ── Build ─────────────────────────────────────────
slide_title()
slide_contents()
slide_goal()
slide_system()
slide_data()
slide_preprocess()
slide_model_arch()
slide_transfer()
slide_training()
slide_confusion()
slide_result()
slide_future()
slide_end()

prs.save(OUTPUT)
print(f"PPT 생성 완료: {OUTPUT}  ({prs.slides.__len__()}슬라이드)")
